# benford/utils.py
import os
import re
import math
import json
import csv
import tempfile
import subprocess
import logging
from collections import Counter
from typing import List, Tuple, Dict, Optional

# —— 日志与调试开关 ——
LOGGER = logging.getLogger(__name__)
DEBUG_OCR = os.environ.get("BENFORD_DEBUG_OCR", "1").lower() not in {"0", "false", "off", ""}


def _log(msg: str):
    if DEBUG_OCR:
        print(f"[OCR] {msg}")
        try:
            LOGGER.info(msg)
        except Exception:
            pass


# —— 文件类型专用库（按需安装） ——
from docx import Document  # .docx
from openpyxl import load_workbook  # .xlsx
from pptx import Presentation  # .pptx
import fitz  # PyMuPDF
from PIL import Image, ImageSequence, ImageEnhance, ImageOps, ImageFilter
import pandas as pd
from bs4 import BeautifulSoup
from lxml import etree
from striprtf.striprtf import rtf_to_text
import filetype

# —— OCR（PaddleOCR，懒加载） ——
from paddleocr import PaddleOCR

_OCR_CACHE: Dict[str, PaddleOCR] = {}


def _get_ocr(lang: str) -> PaddleOCR:
    """
    获取指定语言的 PaddleOCR 实例（PaddleOCR 3.1 适配）。
    说明：
      - 3.1 推荐使用 predict()；无需在 predict() 里再传 cls/angle 参数
      - 关闭文档方向/矫正/文本行方向模块，与你给的示例一致
    """
    if lang not in _OCR_CACHE:
        _log(f"加载 PaddleOCR 模型：lang={lang}")
        _OCR_CACHE[lang] = PaddleOCR(
            lang=lang,
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
        )
    return _OCR_CACHE[lang]


# =========================
# 文本抽取：按类型的优先策略
# =========================
def _read_txt_like(file_path: str, encoding: Optional[str] = None) -> str:
    encs = [encoding, 'utf-8', 'utf-8-sig', 'gb18030', 'latin-1']
    for enc in [e for e in encs if e]:
        try:
            with open(file_path, 'r', encoding=enc, errors='ignore') as f:
                t = f.read()
                _log(f"读取文本文件({enc})：len={len(t)} path={file_path}")
                return t
        except Exception as e:
            _log(f"文本读取失败({enc})：{e}")
            continue
    return ""


def _extract_text_txt(file_path: str) -> str:
    return _read_txt_like(file_path)


def _extract_text_csv(file_path: str) -> str:
    try:
        df = pd.read_csv(file_path, dtype=str, encoding='utf-8', engine='python', errors='ignore')
        t = " ".join(df.fillna("").astype(str).values.ravel().tolist())
        _log(f"CSV 解析：rows={len(df)} len={len(t)} path={file_path}")
        return t
    except Exception as e:
        _log(f"CSV pandas 读取失败：{e}，尝试纯文本模式")
        return _read_txt_like(file_path)


def _extract_text_json(file_path: str) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
        t = json.dumps(data, ensure_ascii=False)
        _log(f"JSON 解析：len={len(t)} path={file_path}")
        return t
    except Exception as e:
        _log(f"JSON 解析失败：{e}，尝试纯文本模式")
        return _read_txt_like(file_path)


def _extract_text_xml(file_path: str) -> str:
    try:
        parser = etree.XMLParser(recover=True)
        root = etree.parse(file_path, parser=parser)
        t = " ".join(root.xpath('//text()'))
        _log(f"XML 解析：len={len(t)} path={file_path}")
        return t
    except Exception as e:
        _log(f"XML 解析失败：{e}，尝试纯文本模式")
        return _read_txt_like(file_path)


def _extract_text_html(file_path: str) -> str:
    try:
        html = _read_txt_like(file_path)
        soup = BeautifulSoup(html, 'html.parser')
        t = soup.get_text(separator=' ')
        _log(f"HTML 提取：len={len(t)} path={file_path}")
        return t
    except Exception as e:
        _log(f"HTML 提取失败：{e}")
        return ""


def _extract_text_md(file_path: str) -> str:
    return _read_txt_like(file_path)


def _extract_text_rtf(file_path: str) -> str:
    try:
        raw = _read_txt_like(file_path)
        t = rtf_to_text(raw)
        _log(f"RTF 提取：len={len(t)} path={file_path}")
        return t
    except Exception as e:
        _log(f"RTF 提取失败：{e}")
        return ""


def _extract_text_docx(file_path: str) -> str:
    doc = Document(file_path)
    parts = [p.text for p in doc.paragraphs]
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                parts.append(cell.text)
    t = "\n".join(parts)
    _log(f"DOCX 提取：paras={len(doc.paragraphs)} len={len(t)} path={file_path}")
    return t


def _extract_text_xlsx(file_path: str) -> str:
    wb = load_workbook(filename=file_path, read_only=True, data_only=True)
    buff = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for v in row:
                if v is not None:
                    buff.append(str(v))
    t = " ".join(buff)
    _log(f"XLSX 提取：sheets={len(wb.worksheets)} tokens={len(buff)} len={len(t)} path={file_path}")
    return t


def _extract_text_xls(file_path: str) -> str:
    try:
        df = pd.read_excel(file_path, dtype=str, engine='xlrd')
        t = " ".join(df.fillna("").astype(str).values.ravel().tolist())
        _log(f"XLS 解析(pandas/xlrd)：rows={len(df)} len={len(t)} path={file_path}")
        return t
    except Exception as e:
        _log(f"XLS 读取失败：{e}（稍后走 LibreOffice→PDF→OCR 兜底）")
        return ""


def _extract_text_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    t = "\n".join(parts)
    _log(f"PPTX 提取：slides={len(prs.slides)} len={len(t)} path={file_path}")
    return t


# ========= 图片 / PDF =========
def _preprocess_image_for_ocr(src_path: str) -> str:
    im = Image.open(src_path).convert("RGB")
    w, h = im.size
    scale = 1.0
    max_side = max(w, h)
    if max_side < 900:
        scale = 900 / max_side
    elif max_side > 2200:
        scale = 2200 / max_side
    if abs(scale - 1.0) > 0.01:
        im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
    im = ImageOps.grayscale(im)
    im = ImageEnhance.Contrast(im).enhance(1.6)
    im = im.filter(ImageFilter.MedianFilter(3))
    im = im.point(lambda x: 255 if x > 200 else (0 if x < 80 else x))
    out_dir = tempfile.mkdtemp(prefix="ocr_pre_")
    out_path = os.path.join(out_dir, "pre.png")
    im.convert("L").save(out_path, "PNG")
    _log(f"图片预处理：src=({w}x{h}) scale={round(scale, 2)} → {out_path}")
    return out_path


def _extract_texts_from_paddle_predict(result) -> List[str]:
    """
    解析 PaddleOCR 3.1 的 predict() 返回结果：
      for res in result:
          res.json['res']['rec_texts'] -> list[str]
    """
    texts: List[str] = []
    if not result:
        return texts
    try:
        for res in result:
            try:
                j = getattr(res, "json", None) or {}
                jres = j.get("res") or {}
                rec_texts = jres.get("rec_texts") or []
                for t in rec_texts:
                    if t:
                        texts.append(str(t))
            except Exception:
                continue
    except Exception:
        pass
    return texts


def _extract_texts_from_ocr_result(res) -> List[str]:
    """
    兼容旧版 PaddleOCR 的 ocr() 返回结构，尽可能抽取 text 列表。
    常见结构：
      - [ [ [box, (text, score) ], ... ] ]  # 单张图片
      - [ [box, (text, score) ], ... ]      # 已经去掉外层的情况
      - [{'text': 'xxx', 'score': ...}, ...]
    """
    texts: List[str] = []
    if not res:
        return texts
    data = res
    if isinstance(res, list) and len(res) == 1 and isinstance(res[0], list):
        data = res[0]
    for item in data:
        try:
            # 结构一： [box, (text, score)]
            if isinstance(item, (list, tuple)) and len(item) >= 2:
                maybe = item[1]
                if isinstance(maybe, (list, tuple)) and len(maybe) >= 1:
                    texts.append(str(maybe[0]))
                    continue
            # 结构二： dict 带 text
            if isinstance(item, dict) and "text" in item:
                texts.append(str(item["text"]))
                continue
        except Exception:
            continue
    return texts


def _ocr_single_image(path: str) -> str:
    pre = _preprocess_image_for_ocr(path)
    text = ""

    # pass1: en（先识别数字英文）
    try:
        ocr_en = _get_ocr("en")
        # —— 优先使用 PaddleOCR 3.1 的 predict()
        res = ocr_en.predict(pre)
        texts = _extract_texts_from_paddle_predict(res)
        _log(f"OCR(en/predict)：boxes={len(texts)} file={path}")
        if not texts:
            # 兜底：尝试旧接口 ocr()
            res_old = ocr_en.ocr(pre)
            texts = _extract_texts_from_ocr_result(res_old)
            _log(f"OCR(en/ocr) 兜底：boxes={len(texts)} file={path}")
        if texts:
            _log(f"OCR(en) sample: {texts[:5]}")
            text = "\n".join(texts)
    except Exception as e:
        _log(f"OCR(en) 异常：{e}")

    # pass2: ch（如 en 无结果，再试中文模型）
    if not text.strip():
        try:
            ocr_ch = _get_ocr("ch")
            res = ocr_ch.predict(pre)
            texts = _extract_texts_from_paddle_predict(res)
            _log(f"OCR(ch/predict)：boxes={len(texts)} file={path}")
            if not texts:
                res_old = ocr_ch.ocr(pre)
                texts = _extract_texts_from_ocr_result(res_old)
                _log(f"OCR(ch/ocr) 兜底：boxes={len(texts)} file={path}")
            if texts:
                _log(f"OCR(ch) sample: {texts[:5]}")
                text = "\n".join(texts)
        except Exception as e:
            _log(f"OCR(ch) 异常：{e}")

    if not text.strip():
        _log(f"OCR 结果为空：{path}")
    return text


def _ocr_images(image_paths: List[str]) -> str:
    parts = []
    for idx, p in enumerate(image_paths, start=1):
        _log(f"开始 OCR 第 {idx}/{len(image_paths)} 张：{p}")
        try:
            t = _ocr_single_image(p)
            if t.strip():
                parts.append(t)
        except Exception as e:
            _log(f"OCR 单图失败：{p}，err={e}")
            continue
    joined = "\n".join(parts)
    _log(f"OCR 多图汇总：images={len(image_paths)} len={len(joined)}")
    return joined


def _convert_pdf_to_images(pdf_path: str, dpi: int = 220) -> List[str]:
    out_dir = tempfile.mkdtemp(prefix="pdf2img_")
    img_paths = []
    with fitz.open(pdf_path) as doc:
        _log(f"PDF 渲染为图片：pages={doc.page_count} dpi={dpi} file={pdf_path}")
        for i, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=dpi)
            p = os.path.join(out_dir, f"page_{i:04d}.png")
            pix.save(p)
            img_paths.append(p)
    _log(f"PDF 渲染完成：images={len(img_paths)} out_dir={out_dir}")
    return img_paths


def _extract_text_pdf(file_path: str) -> str:
    direct = []
    with fitz.open(file_path) as doc:
        for pg in doc:
            try:
                t = pg.get_text() or ""
                if t.strip():
                    direct.append(t)
            except Exception as e:
                _log(f"PDF get_text 异常：{e}")
                continue
    joined = "\n".join(direct).strip()
    _log(f"PDF 直接文本长度：len={len(joined)} file={file_path}")
    if joined:
        return joined
    _log(f"PDF 无文字层，走 OCR：{file_path}")
    imgs = _convert_pdf_to_images(file_path, dpi=220)
    return _ocr_images(imgs)


def _extract_text_image(file_path: str) -> str:
    _log(f"图片走 OCR：{file_path}")
    return _ocr_single_image(file_path)


# =========================
# 兜底：文件转图片 + OCR
# =========================
def _office_to_pdf_via_libreoffice(file_path: str) -> Optional[str]:
    out_dir = tempfile.mkdtemp(prefix="lo_pdf_")
    cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, file_path]
    _log(f"调用 LibreOffice 转换：{' '.join(cmd)}")
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        base = os.path.splitext(os.path.basename(file_path))[0]
        pdf_path = os.path.join(out_dir, base + ".pdf")
        if os.path.exists(pdf_path):
            _log(f"LibreOffice 转换成功：{pdf_path}")
            return pdf_path
        _log("LibreOffice 转换后未找到 PDF 文件")
        return None
    except FileNotFoundError:
        _log("soffice 未安装或不在 PATH，无法进行 Office→PDF 转换")
        return None
    except Exception as e:
        _log(f"LibreOffice 转换失败：{e}")
        return None


def _gif_first_frame_to_png(gif_path: str) -> str:
    out_dir = tempfile.mkdtemp(prefix="gif2png_")
    out = os.path.join(out_dir, "frame_0001.png")
    with Image.open(gif_path) as im:
        for frame in ImageSequence.Iterator(im):
            frame.convert("RGB").save(out, format="PNG")
            break
    _log(f"GIF 首帧导出：{out}")
    return out


def _fallback_to_ocr(file_path: str, ext: str) -> str:
    ext = (ext or "").lower()
    _log(f"进入兜底 OCR，ext={ext} file={file_path}")
    if ext == ".pdf":
        try:
            imgs = _convert_pdf_to_images(file_path, dpi=220)
            return _ocr_images(imgs)
        except Exception as e:
            _log(f"PDF 兜底转换失败：{e}")
            return ""
    office_like = {".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".rtf", ".odt", ".ods", ".odp"}
    if ext in office_like or ext == "":
        pdf = _office_to_pdf_via_libreoffice(file_path)
        if pdf:
            try:
                imgs = _convert_pdf_to_images(pdf, dpi=220)
                return _ocr_images(imgs)
            except Exception as e:
                _log(f"Office→PDF 后渲染/OCR 失败：{e}")
    image_like = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
    if ext in image_like:
        return _ocr_images([file_path])
    if ext == ".gif":
        try:
            one = _gif_first_frame_to_png(file_path)
            return _ocr_images([one])
        except Exception as e:
            _log(f"GIF 首帧导出失败：{e}")
            return ""
    pdf = _office_to_pdf_via_libreoffice(file_path)
    if pdf:
        try:
            imgs = _convert_pdf_to_images(pdf, dpi=220)
            return _ocr_images(imgs)
        except Exception as e:
            _log(f"未知类型→PDF 后渲染/OCR 失败：{e}")
    try:
        with Image.open(file_path) as im:
            tmp_dir = tempfile.mkdtemp(prefix="any2png_")
            out = os.path.join(tmp_dir, "page.png")
            im.convert("RGB").save(out, "PNG")
            _log(f"未知类型当图片处理：{out}")
            return _ocr_images([out])
    except Exception as e:
        _log(f"未知类型当图片处理失败：{e}")
        return ""


# =========================
# 统一入口：获取纯文本 & 提取数字
# =========================
def get_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if not ext:
        try:
            kind = filetype.guess(file_path)
            if kind and kind.EXTENSION:
                ext = "." + kind.EXTENSION.lower()
                _log(f"通过文件签名猜测扩展名：{ext} file={file_path}")
        except Exception as e:
            _log(f"文件签名识别失败：{e}")
    _log(f"开始解析：ext={ext} file={file_path}")
    text = ""
    try:
        if ext in {".txt", ".log"}:
            text = _extract_text_txt(file_path)
        elif ext in {".csv"}:
            text = _extract_text_csv(file_path)
        elif ext in {".json"}:
            text = _extract_text_json(file_path)
        elif ext in {".xml"}:
            text = _extract_text_xml(file_path)
        elif ext in {".html", ".htm"}:
            text = _extract_text_html(file_path)
        elif ext in {".md", ".markdown"}:
            text = _extract_text_md(file_path)
        elif ext in {".rtf"}:
            text = _extract_text_rtf(file_path)
        elif ext in {".docx"}:
            text = _extract_text_docx(file_path)
        elif ext in {".xlsx"}:
            text = _extract_text_xlsx(file_path)
        elif ext in {".xls"}:
            text = _extract_text_xls(file_path)
        elif ext in {".pptx"}:
            text = _extract_text_pptx(file_path)
        elif ext in {".pdf"}:
            text = _extract_text_pdf(file_path)
        elif ext in {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}:
            text = _extract_text_image(file_path)
        elif ext in {".gif"}:
            text = _fallback_to_ocr(file_path, ext)
        else:
            _log("未知扩展名，直接进入兜底 OCR")
            text = ""
    except Exception as e:
        _log(f"直接解析异常：{e}")
        text = ""
    if not text.strip():
        _log("直接解析为空或失败，进入兜底 OCR")
        text = _fallback_to_ocr(file_path, ext)
    _log(f"最终文本长度：len={len(text)} file={file_path}")
    return text or ""


def extract_numbers_from_file(file_path: str) -> List[str]:
    text = get_text_from_file(file_path)
    raw_nums = re.findall(r"\d+", text or "")
    cleaned = []
    for s in raw_nums:
        s2 = s.lstrip("0")
        if not s2:
            continue
        cleaned.append(s2)
    if DEBUG_OCR:
        sample = ", ".join(cleaned[:20])
        _log(f"数字提取：count={len(cleaned)} sample=[{sample}] file={file_path}")
    return cleaned


def extract_numbers_with_context_from_file(
    file_path: str, context_chars: int = 20
) -> List[Dict[str, str]]:
    """从文件中提取数字及其上下文。

    返回列表，每项包含 {"number": 数字字符串, "context": 上下文字符串}
    """
    text = get_text_from_file(file_path)
    results: List[Dict[str, str]] = []
    for m in re.finditer(r"\d+", text or ""):
        raw = m.group(0)
        s = raw.lstrip("0")
        if not s:
            continue
        start, end = m.span()
        ctx = text[max(0, start - context_chars): min(len(text), end + context_chars)]
        ctx = ctx.replace("\n", " ")
        results.append({"number": s, "context": ctx})
    if DEBUG_OCR:
        sample = ", ".join(r["number"] for r in results[:20])
        _log(
            f"数字提取（含上下文）：count={len(results)} sample=[{sample}] file={file_path}"
        )
    return results


# =========================
# Benford 分析
# =========================
def analyze_benford(numbers: List[str]) -> Tuple[Optional[Dict], Optional[Dict], str]:
    if not numbers:
        return None, None, "未提取到任何数字。"
    first_digits = [s[0] for s in numbers if s and s[0].isdigit() and s[0] != "0"]
    total = len(first_digits)
    if total == 0:
        return None, None, "未提取到有效的数字。"
    counts = Counter(first_digits)
    actual = {d: counts.get(str(d), 0) / total * 100 for d in range(1, 10)}
    expected = {d: math.log10(1 + 1 / d) * 100 for d in range(1, 10)}
    max_dev = max(abs(actual[d] - expected[d]) for d in range(1, 10))
    conclusion = ("数值分布偏离本福特定律，存在异常的可能性。"
                  if max_dev > 5.0 else
                  "数值分布与本福特定律总体一致。")
    return actual, expected, conclusion
